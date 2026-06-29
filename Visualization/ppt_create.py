from os import listdir, path
from pptx import Presentation
from pptx.util import Inches, Pt

def create_ppt_from_images(input_dir, outF, ppt_name, title=None):
    # Create a presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    try:
        image_files = sorted([f for f in listdir(input_dir) if f.endswith(('.png', '.jpg', '.jpeg', '.bmp', '.gif'))])
        print(f'Found {len(image_files)} image files in directory')
    except Exception as e:
        print(f"Error accessing input directory: {str(e)}")
        return

    # Process each image file
    for k, image_file in enumerate(image_files, 1):
        image_path = path.join(input_dir, image_file)
        
        print(f'Processing image {k}/{len(image_files)}: {image_file}')
        try:
            # Add a blank slide
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            if title:
                # Add title with filename (without extension)
                title_box = slide.shapes.add_textbox(Inches(0.25), Inches(0.5), Inches(9.5), Inches(0.8))
                text_frame = title_box.text_frame
                text_frame.text = path.splitext(image_file)[0]  # Remove file extension
                
                # Format the title
                paragraph = text_frame.paragraphs[0]
                paragraph.font.size = Pt(32)
                paragraph.font.bold = True
            
            # Add the image to fill most of the slide
            slide.shapes.add_picture(image_path, Inches(0.25), Inches(1.5),
                                    width=Inches(9.5))
            
        except Exception as e:
            print(f'Error processing {image_file}: {str(e)}')
            continue

    # Save the presentation
    ppt_path = path.join(outF, ppt_name+'.pptx')
    prs.save(ppt_path)
    print(f"Presentation saved at: {ppt_path}")

    print(f'\n{"="*60}')
    print(f'Total images processed: {k}/{len(image_files)}')
    print(f'PowerPoint saved to: {ppt_path}')


if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser(description='Create a PowerPoint presentation from images in a directory.')
    parser.add_argument('--directory', type=str, required=False, help='Directory containing image files.')
    parser.add_argument('--output', type=str, required=False, help='Directory to save the PowerPoint presentation.')
    parser.add_argument('--ppt_name', type=str, default='images_presentation.pptx', help='Name of the output PowerPoint file.')
    parser.add_argument('--title', type=str, default=None, help='Filename as title for each slide.')
    args = parser.parse_args()  

    if args.directory and args.output:
        create_ppt_from_images(args.directory, args.output, args.ppt_name, args.title)

# # Define directories
# input_dir = r"C:\Users\lilif\OneDrive\Desktop\Dropbox\Phd\Pipeline testing\Images\3Im"
# outF = r"C:\Users\lilif\OneDrive\Desktop\Dropbox\Phd\Pipeline testing\Images"

# create_ppt_from_images(input_dir, outF, 'images_presentation.pptx', title=None)